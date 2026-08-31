"""
PathFinder Nexus — PowerPoint Presentation Generator
Creates a 16:9 widescreen, dark-themed, 15-slide PowerPoint deck for college evaluation, viva, or investor pitch.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_COLOR = RGBColor(2, 6, 23)        # Slate-950
    CARD_BG = RGBColor(15, 23, 42)       # Slate-900
    BORDER_COLOR = RGBColor(30, 41, 59)  # Slate-800
    CYAN = RGBColor(6, 182, 212)         # #06b6d4
    EMERALD = RGBColor(16, 185, 129)     # #10b981
    TEAL = RGBColor(20, 184, 166)        # #14b8a6
    WHITE = RGBColor(248, 250, 252)      # #f8fafc
    SLATE_LIGHT = RGBColor(203, 213, 225)# #cbd5e1
    SLATE_MUTED = RGBColor(148, 163, 184)# #94a3b8
    ACCENT_ROSE = RGBColor(244, 63, 94)  # #f43f5e
    ACCENT_AMBER = RGBColor(245, 158, 11)# #f59e0b

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text, tag_color=CYAN):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = tag_color

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    # ==========================================================
    # SLIDE 1: Title Slide
    # ==========================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Center hero card
    card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CYAN
    card.line.width = Pt(2)

    # Title box
    t_box = s1.shapes.add_textbox(Inches(1.8), Inches(1.6), Inches(9.7), Inches(4.3))
    tf = t_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "FULL-STACK CAPSTONE PROJECT DEFENSE"
    p0.alignment = PP_ALIGN.CENTER
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = CYAN

    p1 = tf.add_paragraph()
    p1.text = "PathFinder Nexus"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "Autonomous, Dependency-Aware & Continuously Adaptive Learning Navigation Platform"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(18)
    p2.font.color.rgb = SLATE_LIGHT
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "FastAPI (Python 3.12) • PostgreSQL 16 & pgvector • React 18 SPA • Google Gemini LLM"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = EMERALD
    p3.space_before = Pt(25)

    p4 = tf.add_paragraph()
    p4.text = "162 Backend Tests Passing • 100% Deterministic DAG Engine & Bayesian Evidence Fusion"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.size = Pt(12)
    p4.font.color.rgb = SLATE_MUTED
    p4.space_before = Pt(8)

    # ==========================================================
    # SLIDE 2: Problem Statement
    # ==========================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "01 / The Core Problem", "Why Modern Online Education Fails Learners", ACCENT_ROSE)

    problems = [
        ("Linear, Static Playlists", "Courses assume identical starting points, forcing experienced learners through basics and abandoning beginners when advanced prerequisites are skipped.", ACCENT_ROSE),
        ("Hallucinating AI Chatbots", "Generic LLMs invent fake URLs, fabricate invalid course sequences, and lack relational database grounding and dependency validation.", ACCENT_AMBER),
        ("Zero Adaptive Feedback Loop", "When learners fail assessments, traditional platforms do not dynamically restructure roadmaps or deliver prerequisite interventions.", CYAN),
    ]

    for i, (title, desc, color) in enumerate(problems):
        x = Inches(0.8 + i * 4.0)
        c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s2.shapes.add_textbox(x + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p_num = tf.paragraphs[0]
        p_num.text = f"0{i+1}"
        p_num.font.size = Pt(24)
        p_num.font.bold = True
        p_num.font.color.rgb = color

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_before = Pt(14)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = SLATE_MUTED
        p_d.space_before = Pt(12)

    # ==========================================================
    # SLIDE 3: Architectural Creed
    # ==========================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "02 / Architectural Philosophy", "The PathFinder Nexus Architectural Creed", CYAN)

    quote_card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.0))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = CARD_BG
    quote_card.line.color.rgb = CYAN
    quote_card.line.width = Pt(2)

    q_box = s3.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(1.4))
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    qp = q_tf.paragraphs[0]
    qp.text = "“RAG retrieves. LLMs explain. The Database grounds. Deterministic engines decide.”"
    qp.font.size = Pt(24)
    qp.font.bold = True
    qp.font.color.rgb = CYAN
    qp.alignment = PP_ALIGN.CENTER

    pillars = [
        ("RAG (pgvector)", "Retrieves verified educational resources via 1536-dim cosine similarity.", CYAN),
        ("LLM (Gemini)", "Synthesizes clear explanations, goal summaries, and interactive tutor answers.", TEAL),
        ("Database (Postgres)", "Authoritative single source of truth for roles, skills, questions, and roadmaps.", EMERALD),
        ("Deterministic Engines", "Topological Kahn DAG, 6-factor ranking, and Bayesian evidence fusion.", WHITE),
    ]

    for i, (title, desc, color) in enumerate(pillars):
        x = Inches(0.8 + i * 3.0)
        c = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(2.75), Inches(2.5))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = BORDER_COLOR
        c.line.width = Pt(1)

        tb = s3.shapes.add_textbox(x + Inches(0.2), Inches(4.4), Inches(2.35), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = SLATE_MUTED
        pd.space_before = Pt(8)

    # ==========================================================
    # SLIDE 4: System Architecture
    # ==========================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "03 / System Blueprint", "Modular Monolith Architecture & Technology Stack", EMERALD)

    layers = [
        ("Frontend Client (SPA)", [
            "React 18 + TypeScript + Vite",
            "Tailwind CSS Dark Mode & Glassmorphism",
            "Client-Side JWT Auth with auto-refresh",
            "Netlify Edge CDN Deployment"
        ], CYAN),
        ("FastAPI Backend Core", [
            "Python 3.12 Modular Monolith",
            "Clean Architecture (Routers/Services/Repos)",
            "Argon2id Password Security + Correlation IDs",
            "Render Cloud Web Service (Uvicorn)"
        ], EMERALD),
        ("Data & AI Layer", [
            "PostgreSQL 16 with pgvector Extension",
            "22 Normalized Relational Entity Tables",
            "Google Gemini 1.5/2.0 Flash + OpenAI Embeddings",
            "Deterministic Mock Fallback Provider"
        ], TEAL),
    ]

    for i, (title, items, color) in enumerate(layers):
        x = Inches(0.8 + i * 4.0)
        c = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s4.shapes.add_textbox(x + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = color

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = SLATE_LIGHT
            p.space_before = Pt(10)

    # ==========================================================
    # SLIDE 5: Engine 1 — AI Goal Extraction
    # ==========================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "04 / Engine 1", "AI Goal Understanding & Grounded Extraction", CYAN)

    # Left: NL input & extraction
    c_left = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = CYAN

    tb_l = s5.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "Natural Language Input Example:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CYAN

    p = tf_l.add_paragraph()
    p.text = "“I want to become an AI Engineer specializing in LLMs & RAG systems within 12 weeks with 2 hours of daily study.”"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

    p = tf_l.add_paragraph()
    p.text = "Structured Pydantic Extraction Output:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_before = Pt(16)

    p = tf_l.add_paragraph()
    p.text = "• Target Role: AI/ML Engineer (Canonical Match: 95%)\n• Timeline: 12 Weeks\n• Daily Commitment: 2.0 Hours\n• Prior Skills Identified: Python (Intermediate), Git (Basic)"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_LIGHT
    p.space_before = Pt(6)

    # Right: Schema & Grounding
    c_right = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = BORDER_COLOR

    tb_r = s5.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Canonical Role Grounding Pipeline:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEAL

    p = tf_r.add_paragraph()
    p.text = "1. Raw Prompt Sanitization & Delimiter Tagging\n2. Gemini Structured JSON Mode with Pydantic Schema Validation\n3. Database Lookup: Matches against 10 Canonical Roles\n4. Unknown Role Detection: Returns AMBIGUOUS with suggested roles\n5. Zero Hallucination Guarantee: Unverified roles are never persisted"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_LIGHT
    p.space_before = Pt(10)

    # ==========================================================
    # SLIDE 6: Engine 2 — Dynamic Skill Gap
    # ==========================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "05 / Engine 2", "Dynamic Skill Gap & Role Readiness Engine", EMERALD)

    metrics = [
        ("Zero Static Storage", "Skill gaps are calculated dynamically upon query. Eliminates stale cached tables and guarantees synchronization.", CYAN),
        ("Gap Formulation", "Gap = max(Required - LearnerProficiency, 0)\nStrictly non-negative delta for every skill required by role.", EMERALD),
        ("Role Readiness Percentage", "Readiness = (Σ Current_w / Σ Required_w) × 100\nWeighted importance factor ensures critical skills drive readiness.", TEAL),
    ]

    for i, (title, desc, color) in enumerate(metrics):
        x = Inches(0.8 + i * 4.0)
        c = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s6.shapes.add_textbox(x + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = SLATE_LIGHT
        pd.space_before = Pt(14)

    # ==========================================================
    # SLIDE 7: Engine 3 — Topological Roadmap (Kahn's DAG)
    # ==========================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "06 / Engine 3", "Topological Roadmap Engine (Kahn's DAG Algorithm)", TEAL)

    # Left: Theory
    c_left = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = TEAL

    tb_l = s7.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "Prerequisite-Aware Directed Acyclic Graph"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEAL

    p = tf_l.add_paragraph()
    p.text = "• Every skill node defines strict prerequisites in the database\n• In-Degree topological sort ensures prerequisite skills strictly precede dependent topics\n• Cycle detection prevents infinite dependency loops\n• Automatic locking of downstream modules until prerequisites achieve passing mastery"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_LIGHT
    p.space_before = Pt(12)

    # Right: Algorithm Steps
    c_right = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = BORDER_COLOR

    tb_r = s7.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Kahn's Algorithm Execution Workflow:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p = tf_r.add_paragraph()
    p.text = "1. Calculate in-degree for all required skills\n2. Initialize queue with all in-degree == 0 nodes\n3. Pop node u -> append to roadmap sequence\n4. Decrement in-degree for all outgoing neighbors v\n5. If in-degree[v] == 0, push v to queue\n6. Verify all nodes processed; assign milestones & weeks"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_MUTED
    p.space_before = Pt(10)

    # ==========================================================
    # SLIDE 8: Engine 4 — 6-Factor Recommendation
    # ==========================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "07 / Engine 4", "Normalized 6-Factor Explainable Recommendation", CYAN)

    formula_card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.2))
    formula_card.fill.solid()
    formula_card.fill.fore_color.rgb = CARD_BG
    formula_card.line.color.rgb = CYAN

    fb = s8.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.3), Inches(0.9))
    ftf = fb.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Score = 0.30(Gap) + 0.20(Prereq) + 0.15(Goal) + 0.15(Diff) + 0.10(Time) + 0.10(Pref)"
    fp.alignment = PP_ALIGN.CENTER
    fp.font.size = Pt(18)
    fp.font.bold = True
    fp.font.color.rgb = CYAN

    factors = [
        ("30%", "Skill Gap Weight", "Prioritizes resources targeting the largest deficiency.", CYAN),
        ("20%", "Prereq Readiness", "Ensures learner is prepared for resource concepts.", TEAL),
        ("15%", "Goal Alignment", "Matches specific career target role curriculum.", EMERALD),
        ("15%", "Difficulty Match", "Aligns with learner's current experience level.", WHITE),
        ("10%", "Time Investment", "Fits daily and weekly study budget.", ACCENT_AMBER),
        ("10%", "Format Preference", "Balances video, article, and interactive preferences.", ACCENT_ROSE),
    ]

    for i, (pct, title, desc, color) in enumerate(factors):
        x = Inches(0.8 + (i % 3) * 4.0)
        y = Inches(3.3 + (i // 3) * 1.8)
        c = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(1.5))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1)

        tb = s8.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.4), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = f"{pct} — {title}"
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = SLATE_MUTED
        pd.space_before = Pt(4)

    # ==========================================================
    # SLIDE 9: Engine 5 & 6 — Assessment & Evidence Fusion
    # ==========================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "08 / Engine 5 & 6", "Sanitized Assessments & Bayesian Evidence Fusion", EMERALD)

    # Left: Anti-Cheat
    c_left = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = CYAN

    tb_l = s9.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "Anti-Cheat Sanitized Delivery"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN

    p = tf_l.add_paragraph()
    p.text = "• Correct answers and explanation rationales are NEVER sent to the client browser\n• Client receives only question IDs, question text, and option keys\n• Authoritative grading occurs strictly on the server\n• Prevents browser inspection exploits and ensures grading integrity"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_LIGHT
    p.space_before = Pt(12)

    # Right: Evidence Fusion
    c_right = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = EMERALD

    tb_r = s9.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Bayesian-Inspired Evidence Fusion"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_r.add_paragraph()
    p.text = "Mastery Formula:\nP_new = round(0.30 × P_old + 0.70 × AssessmentScore, 2)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

    p = tf_r.add_paragraph()
    p.text = "• Prevents single-test volatility while rewarding demonstrated competence\n• Dynamic mastery thresholding: Mastery (≥80%), Continue (60-79%), Reinforce (<60%)"
    p.font.size = Pt(12)
    p.font.color.rgb = SLATE_MUTED
    p.space_before = Pt(12)

    # ==========================================================
    # SLIDE 10: Closed-Loop Adaptive Interventions
    # ==========================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "09 / Adaptive Core", "Closed-Loop Adaptive Interventions & Next Best Action", ACCENT_AMBER)

    tiers = [
        ("Score < 40% — Critical Gap", "Foundational Lock", "Locks dependent roadmap items, schedules foundational refresher resources, and updates Next Best Action to reinforce fundamentals.", ACCENT_ROSE),
        ("40% ≤ Score < 80% — Partial", "Targeted Reinforcement", "Generates targeted practice resources while keeping milestone active for revision.", ACCENT_AMBER),
        ("Score ≥ 80% — Mastery", "Milestone Unlocked", "Marks milestone completed, updates topological graph, and unlocks downstream advanced milestones.", EMERALD),
    ]

    for i, (tag, title, desc, color) in enumerate(tiers):
        x = Inches(0.8 + i * 4.0)
        c = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s10.shapes.add_textbox(x + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = tag
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = color

        pt2 = tf.add_paragraph()
        pt2.text = title
        pt2.font.size = Pt(18)
        pt2.font.bold = True
        pt2.font.color.rgb = WHITE
        pt2.space_before = Pt(8)

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = SLATE_LIGHT
        pd.space_before = Pt(12)

    # ==========================================================
    # SLIDE 11: Engine 7 — Grounded RAG Assistant
    # ==========================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "10 / Engine 7", "Grounded RAG Knowledge Assistant & Citation Safety", CYAN)

    rag_cards = [
        ("pgvector Semantic Search", [
            "1536-dimensional embeddings for all curated catalog resources",
            "Cosine similarity threshold cutoff (≥ 0.50)",
            "Top-K semantic context retrieval directly from PostgreSQL"
        ], CYAN),
        ("Anti-Hallucination Guardrails", [
            "Retrieved database content wrapped in strict XML security delimiters",
            "LLM must cite authentic resource URLs present in retrieved context",
            "Zero-context fallback: refuses to fabricate imaginary resources"
        ], EMERALD),
    ]

    for i, (title, points, color) in enumerate(rag_cards):
        x = Inches(0.8 + i * 6.0)
        c = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(5.6), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s11.shapes.add_textbox(x + Inches(0.3), Inches(2.1), Inches(5.0), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = color

        for p_text in points:
            p = tf.add_paragraph()
            p.text = f"• {p_text}"
            p.font.size = Pt(12)
            p.font.color.rgb = SLATE_LIGHT
            p.space_before = Pt(12)

    # ==========================================================
    # SLIDE 12: Database Schema (22 Core Tables)
    # ==========================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_header(s12, "11 / Data Foundation", "Relational Database Schema (22 Core Tables)", EMERALD)

    db_groups = [
        ("Auth & Users", "users\nlearner_profiles\nlearner_skills\nuser_preferences", CYAN),
        ("Catalog & Graph", "roles\nskills\nskill_prerequisites\nrole_skills", TEAL),
        ("Roadmaps & Learning", "roadmaps\nroadmap_items\nresources\nprojects", EMERALD),
        ("Assessments & RAG", "assessments\nassessment_questions\nsubmissions\nchat_messages", WHITE),
    ]

    for i, (title, tables, color) in enumerate(db_groups):
        x = Inches(0.8 + i * 3.0)
        c = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1)

        tb = s12.shapes.add_textbox(x + Inches(0.2), Inches(2.1), Inches(2.35), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = color

        pd = tf.add_paragraph()
        pd.text = tables
        pd.font.size = Pt(12)
        pd.font.color.rgb = SLATE_LIGHT
        pd.space_before = Pt(14)

    # ==========================================================
    # SLIDE 13: Testing & Quality Metrics
    # ==========================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13)
    add_header(s13, "12 / Verification", "Automated Testing & Benchmark Results", TEAL)

    test_cards = [
        ("162 Tests", "Backend Pytest Suite", "100% passing across auth, DAG roadmaps, recommendations, evidence fusion, and security.", EMERALD),
        ("0 Errors", "TypeScript Strict Build", "Zero TypeScript compilation errors with optimized Vite production bundle.", CYAN),
        ("100% Safe", "Deterministic Fallback", "System operates smoothly even when cloud AI rate limits hit or in offline mode.", TEAL),
    ]

    for i, (stat, title, desc, color) in enumerate(test_cards):
        x = Inches(0.8 + i * 4.0)
        c = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s13.shapes.add_textbox(x + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p_stat = tf.paragraphs[0]
        p_stat.text = stat
        p_stat.font.size = Pt(36)
        p_stat.font.bold = True
        p_stat.font.color.rgb = color

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = WHITE
        p_t.space_before = Pt(10)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SLATE_MUTED
        p_d.space_before = Pt(12)

    # ==========================================================
    # SLIDE 14: Cloud Deployment Architecture
    # ==========================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_background(s14)
    add_header(s14, "13 / Deployment", "Production Cloud Deployment Architecture", CYAN)

    deploy_tiers = [
        ("Frontend: Netlify", [
            "React 18 SPA on Global Edge CDN",
            "Automated netlify.toml headers & caching",
            "_redirects file for SPA client-side deep linking",
            "Command: npm run build"
        ], CYAN),
        ("Backend: Render", [
            "FastAPI Web Service running Uvicorn",
            "render.yaml blueprint configuration",
            "Dynamic CORS regex matching Netlify previews",
            "Health Check Probes: /health & /health/ready"
        ], EMERALD),
        ("Database: PostgreSQL", [
            "Cloud PostgreSQL 16 Instance",
            "Automated migrations via deploy_init.py",
            "Idempotent catalog seeding on startup",
            "pgvector similarity indexing"
        ], TEAL),
    ]

    for i, (title, items, color) in enumerate(deploy_tiers):
        x = Inches(0.8 + i * 4.0)
        c = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = color
        c.line.width = Pt(1.5)

        tb = s14.shapes.add_textbox(x + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = color

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = SLATE_LIGHT
            p.space_before = Pt(10)

    # ==========================================================
    # SLIDE 15: Conclusion
    # ==========================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_background(s15)

    card = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = EMERALD
    card.line.width = Pt(2)

    tb = s15.shapes.add_textbox(Inches(1.8), Inches(1.6), Inches(9.7), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "CONCLUSION & LIVE DEMO"
    p0.alignment = PP_ALIGN.CENTER
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = EMERALD

    p1 = tf.add_paragraph()
    p1.text = "PathFinder Nexus"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "“Transforming natural-language ambition into deterministic, adaptive mastery.”"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(16)
    p2.font.italic = True
    p2.font.color.rgb = CYAN
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "✓ 7 Fully Integrated Deterministic Engines  |  ✓ 162 Passed Tests  |  ✓ Deployment-Ready"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = SLATE_LIGHT
    p3.space_before = Pt(24)

    p4 = tf.add_paragraph()
    p4.text = "Thank you! We are ready for live demonstration & questions."
    p4.alignment = PP_ALIGN.CENTER
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEAL
    p4.space_before = Pt(16)

    # Save file
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "PathFinder_Nexus_Project_Presentation.pptx")
    prs.save(output_path)
    print(f"[SUCCESS] PowerPoint Presentation successfully created: {output_path}")

if __name__ == "__main__":
    create_presentation()
